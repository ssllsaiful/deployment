import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Theme Colors (Deep Corporate Blue & Teal Accent)
DARK_BG = RGBColor(10, 25, 47)       # Navy
LIGHT_BG = RGBColor(245, 247, 250)    # Soft Grayish White
TEXT_DARK = RGBColor(20, 32, 48)      # Dark Charcoal
TEXT_LIGHT = RGBColor(240, 240, 240)   # Light Gray
TEAL_ACCENT = RGBColor(0, 150, 136)   # Teal/Cyan
CARD_BG = RGBColor(255, 255, 255)     # Card White

def set_slide_background(slide, color):
    """Sets a solid color background for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, main_title, subtitle_text):
    """Adds a dark corporate style title slide."""
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)
    
    # Title Text Box
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    # Title Paragraph
    p1 = tf.paragraphs[0]
    p1.text = main_title
    p1.font.bold = True
    p1.font.size = Pt(44)
    p1.font.color.rgb = TEAL_ACCENT
    p1.font.name = "Arial"
    p1.alignment = PP_ALIGN.LEFT
    
    # Space
    p_space = tf.add_paragraph()
    p_space.text = ""
    p_space.font.size = Pt(14)
    
    # Subtitle Paragraph
    p2 = tf.add_paragraph()
    p2.text = subtitle_text
    p2.font.size = Pt(20)
    p2.font.color.rgb = TEXT_LIGHT
    p2.font.name = "Arial"
    p2.alignment = PP_ALIGN.LEFT

def add_content_slide(prs, step_label, slide_title, details_list):
    """Adds a clean content slide with a side panel indicator and formatted bullet points."""
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, LIGHT_BG)
    
    # Decorative Top Accent Bar
    accent_bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.333), Inches(0.15) # Rectangular shape
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = TEAL_ACCENT
    accent_bar.line.color.rgb = TEAL_ACCENT

    # Step Label Text Box (Top-Left)
    stepBox = slide.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.33), Inches(0.5))
    stf = stepBox.text_frame
    stf.word_wrap = True
    sp = stf.paragraphs[0]
    sp.text = step_label.upper()
    sp.font.bold = True
    sp.font.size = Pt(14)
    sp.font.color.rgb = TEAL_ACCENT
    sp.font.name = "Arial"
    
    # Title Text Box (Below Step Label)
    titleBox = slide.shapes.add_textbox(Inches(1.0), Inches(0.9), Inches(11.33), Inches(0.8))
    ttf = titleBox.text_frame
    ttf.word_wrap = True
    tp = ttf.paragraphs[0]
    tp.text = slide_title
    tp.font.bold = True
    tp.font.size = Pt(28)
    tp.font.color.rgb = TEXT_DARK
    tp.font.name = "Arial"

    # Main Details Box
    detailsBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.33), Inches(5.0))
    dtf = detailsBox.text_frame
    dtf.word_wrap = True
    
    for i, detail in enumerate(details_list):
        if i == 0:
            dp = dtf.paragraphs[0]
        else:
            dp = dtf.add_paragraph()
            
        dp.text = "•  " + detail
        dp.font.size = Pt(16)
        dp.font.color.rgb = TEXT_DARK
        dp.font.name = "Arial"
        dp.space_after = Pt(14)
        
        # Check if line contains command codes, formats differently
        if "$" in detail or "sudo" in detail or "npm" in detail or "pm2" in detail:
            dp.font.name = "Courier New"
            dp.font.size = Pt(14)
            dp.font.color.rgb = RGBColor(30, 41, 59)

def generate_frontend():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Title
    add_title_slide(
        prs, 
        "Next.js Frontend Deployment", 
        "Step-by-Step Production Guide on Ubuntu Server via PM2 & Nginx"
    )
    
    # Prerequisites
    add_content_slide(
        prs, 
        "Prerequisites", 
        "Required Credentials & Access Setup",
        [
            "Ensure you have the server public IP address and the SSH login credential (private .pem key or root password).",
            "Point your domain names (DNS A records) directly to your server's IP address.",
            "Verify you have access to pull the code from the Git repository on the server.",
            "Gather any frontend environment variables (like API endpoint URLs) to copy into your .env configuration file."
        ]
    )

    # Step 1
    add_content_slide(
        prs,
        "Step 01",
        "Establish SSH Connection & Update Server Packages",
        [
            "Open your local terminal (macOS/Linux) or PowerShell (Windows).",
            "Connect to the server using SSH:",
            "  $ ssh -i /path/to/key.pem ubuntu@your_server_ip",
            "Update system libraries to ensure security patches are applied:",
            "  $ sudo apt update && sudo apt upgrade -y"
        ]
    )

    # Step 2
    add_content_slide(
        prs,
        "Step 02",
        "Install Node.js Version Manager (NVM) & Node.js",
        [
            "Install NVM on the server to isolate and manage Node.js runtimes:",
            "  $ curl -o- https://raw.githubusercontent.com/nvm-sh/nvm/v0.39.7/install.sh | bash",
            "Reload environment settings to enable the nvm command:",
            "  $ source ~/.bashrc   (or export NVM variables directly in terminal)",
            "Install Node.js version 20 (LTS):",
            "  $ nvm install 20",
            "Confirm node and npm installations using 'node -v' and 'npm -v'."
        ]
    )

    # Step 3
    add_content_slide(
        prs,
        "Step 03",
        "Create Directory, Clone Repo, & Setup Env Variables",
        [
            "Create the deployment folder and adjust ownership rules to avoid permission errors:",
            "  $ sudo mkdir -p /home/ubuntu/frontend && sudo chown -R $USER:$USER /home/ubuntu/frontend",
            "Navigate into the folder and clone your repository:",
            "  $ cd /home/ubuntu/frontend && git clone <your_repository_url> projectname",
            "Navigate to project root, create a production environment variables file:",
            "  $ cd projectname && nano .env",
            "Add required configurations (e.g. NODE_ENV=production, NEXT_PUBLIC_API_URL=...) and save changes."
        ]
    )

    # Step 4
    add_content_slide(
        prs,
        "Step 04",
        "Install Project Dependencies & Build Next.js",
        [
            "Install local dependencies including devDependencies needed to build the Next.js code:",
            "  $ npm install",
            "Build the Next.js production bundle:",
            "  $ npm run build",
            "This compiles static pages, components, CSS styles, and templates inside the hidden '.next' directory.",
            "Verify build is completed successfully without errors."
        ]
    )

    # Step 5
    add_content_slide(
        prs,
        "Step 05",
        "Configure PM2 Process Manager & Ecosystem File",
        [
            "Install PM2 globally on the server to monitor processes:",
            "  $ sudo npm install pm2 -g",
            "Create a configuration file 'ecosystem.config.js' at the root of the project:",
            "  $ nano ecosystem.config.js",
            "Add app configurations, specifying name as 'frontend', script as 'npm', and args as 'start'.",
            "Start the frontend using the ecosystem file:   $ pm2 start ecosystem.config.js",
            "Save the configuration and check status:   $ pm2 save && pm2 status"
        ]
    )

    # Step 6
    add_content_slide(
        prs,
        "Step 06",
        "Reverse Proxy Configuration (CloudPanel vs Nginx)",
        [
            "Option A (CloudPanel): Add a Reverse Proxy site for your domain pointing to port 3000, and generate SSL.",
            "Option B (Manual Nginx): Verify or install Nginx:",
            "  $ sudo apt update && sudo apt install nginx -y",
            "Create a file in sites-available:   $ sudo nano /etc/nginx/sites-available/frontend.example.com",
            "Add a server block listening on port 80 proxying requests to http://localhost:3000.",
            "Link config to sites-enabled:   $ sudo ln -s /etc/nginx/sites-available/... /etc/nginx/sites-enabled/",
            "Test settings and restart Nginx:   $ sudo nginx -t && sudo systemctl restart nginx"
        ]
    )

    # Step 7
    add_content_slide(
        prs,
        "Step 07",
        "Enable SSL (HTTPS) with Certbot",
        [
            "Install Certbot and the Nginx helper package on the server:",
            "  $ sudo apt install certbot python3-certbot-nginx -y",
            "Generate Let's Encrypt certificates and auto-configure Nginx:",
            "  $ sudo certbot --nginx -d frontend.example.com",
            "Select options to redirect all traffic to HTTPS automatically.",
            "Test automatic renewal configuration:   $ sudo certbot renew --dry-run"
        ]
    )

    # Step 8
    add_content_slide(
        prs,
        "Step 08",
        "Automate Deployment via GitHub Actions",
        [
            "Set up repository secrets in GitHub settings: SSH_HOST, SSH_USER, and SSH_KEY.",
            "Create a workflow file in local project: .github/workflows/deploy.yml",
            "Configure a push action workflow for the 'main' branch.",
            "Define deployment steps: SSH into server, git pull, npm install, npm run build, pm2 restart ecosystem.config.js.",
            "Commit and push deploy.yml to automatically trigger and verify deployment on push."
        ]
    )
    
    output_path = os.path.join("forntend-next.js", "frontend_deployment.pptx")
    prs.save(output_path)
    print(f"Generated: {output_path}")

def generate_backend():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Title
    add_title_slide(
        prs, 
        "Django Backend & PostgreSQL Deployment", 
        "Step-by-Step Production Guide on Ubuntu Server using PM2, Nginx, & GitHub Actions"
    )
    
    # Prerequisites
    add_content_slide(
        prs, 
        "Prerequisites", 
        "Required Parameters & Credentials",
        [
            "Gather the server public IP, SSH user, and private key file (.pem).",
            "Define database info to create: DB Name, DB User, and DB password.",
            "Point your backend API subdomain (DNS A record) to the server IP.",
            "Locate your project settings directory to update ALLOWED_HOSTS and production database configurations."
        ]
    )

    # Step 1
    add_content_slide(
        prs,
        "Step 01",
        "Log in and Install PostgreSQL Database",
        [
            "Establish SSH connection to the server:   $ ssh -i <key> ubuntu@your_server_ip",
            "Update packages and install PostgreSQL plus database compiler dependencies:",
            "  $ sudo apt update && sudo apt install postgresql postgresql-contrib -y"
        ]
    )

    # Step 2
    add_content_slide(
        prs,
        "Step 02",
        "Create Database & User, Grant Privileges",
        [
            "Access the administrative SQL prompt:   $ sudo -i -u postgres psql",
            "Create database:   CREATE DATABASE my_django_db;",
            "Create user:   CREATE USER django_db_user WITH PASSWORD 'SuperSecurePassword123';",
            "Assign all permissions:   GRANT ALL PRIVILEGES ON DATABASE my_django_db TO django_db_user;",
            "On PostgreSQL 15+, switch database and grant schema rights:",
            "  \\c my_django_db",
            "  GRANT ALL ON SCHEMA public TO django_db_user;",
            "Exit database CLI:   \\q"
        ]
    )

    # Step 3
    add_content_slide(
        prs,
        "Step 03",
        "Configure Remote Access for PostgreSQL (Optional)",
        [
            "Edit postgresql.conf to listen on all IP addresses:",
            "  $ sudo nano /etc/postgresql/15/main/postgresql.conf  -> listen_addresses = '*'",
            "Edit pg_hba.conf to whitelist connections:",
            "  $ sudo nano /etc/postgresql/15/main/pg_hba.conf",
            "  Add: host all all 0.0.0.0/0 scram-sha-256 (or limit to frontend IP for security)",
            "Open port 5432 and restart service:",
            "  $ sudo ufw allow 5432/tcp && sudo systemctl restart postgresql"
        ]
    )

    # Step 4
    add_content_slide(
        prs,
        "Step 04",
        "Update Local Django Project Settings",
        [
            "Open the local project on your PC.",
            "Install Gunicorn (WSGI server) and WhiteNoise (for static file assets styling in production):",
            "  $ pip install gunicorn whitenoise",
            "Update settings.py with: DEBUG = False, ALLOWED_HOSTS = ['api.example.com', 'server_ip']",
            "Add 'whitenoise.middleware.WhiteNoiseMiddleware' directly under security middleware.",
            "Configure database parameters to use server-side credentials and set static root pathways."
        ]
    )

    # Step 5
    add_content_slide(
        prs,
        "Step 05",
        "Save Requirements & Push to Git Repository",
        [
            "Export all project requirements:   $ pip freeze > requirements.txt",
            "Confirm gunicorn, whitenoise, and database adapter packages are listed in requirements.txt.",
            "Commit and push your files to remote repo (GitHub/GitLab):",
            "  $ git add . && git commit -m 'Setup production settings' && git push origin main"
        ]
    )

    # Step 6
    add_content_slide(
        prs,
        "Step 06",
        "Setup Server Directory and Pull Code",
        [
            "Create deployment directories on the server:   $ sudo mkdir -p /home/ubuntu/backend",
            "Set user ownership rules:   $ sudo chown -R $USER:$USER /home/ubuntu/backend",
            "Clone repository inside directory:   $ cd /home/ubuntu/backend && git clone <url> repo",
            "Create environment file:   $ cd repo && nano .env",
            "Paste database details, key configurations, debug settings and save the file."
        ]
    )

    # Step 7
    add_content_slide(
        prs,
        "Step 07",
        "Configure Python Virtualenv and Install Packages",
        [
            "Create a python virtual environment:   $ python3 -m venv venv",
            "Activate the environment:   $ source venv/bin/activate",
            "Install project requirement dependencies:   $ pip install --upgrade pip && pip install -r requirements.txt",
            "Run migrations:   $ python manage.py migrate",
            "Generate admin css assets static files:   $ python manage.py collectstatic --no-input"
        ]
    )

    # Step 8
    add_content_slide(
        prs,
        "Step 08",
        "Setup PM2 to Run Gunicorn Server",
        [
            "Create 'ecosystem.config.js' file in backend folder:   $ nano ecosystem.config.js",
            "Define application settings:",
            "  Name: 'backend', Script: 'venv/bin/gunicorn', Args: '--workers 3 --bind 127.0.0.1:8000 myproject.wsgi:application'",
            "Start process manager:   $ pm2 start ecosystem.config.js",
            "Persist running backend application:   $ pm2 save"
        ]
    )

    # Step 9
    add_content_slide(
        prs,
        "Step 09",
        "Reverse Proxy Configuration (CloudPanel vs Nginx)",
        [
            "Option A (CloudPanel): Add a Reverse Proxy targeting port 8000 and configure SSL via dashboard.",
            "Option B (Manual Nginx): Add Nginx configuration file:",
            "  $ sudo nano /etc/nginx/sites-available/api.example.com",
            "Define locations for static files alias paths and proxy pass other routing queries to http://localhost:8000.",
            "Activate site configurations and reboot:   $ sudo ln -s ... && sudo nginx -t && sudo systemctl restart nginx"
        ]
    )

    # Step 10
    add_content_slide(
        prs,
        "Step 10",
        "Configure SSL & Automation via GitHub Actions",
        [
            "Run Certbot for SSL configuration:   $ sudo certbot --nginx -d api.example.com",
            "Configure SSH host, user, and keys under repository secrets in GitHub.",
            "Create local repository workflow file: .github/workflows/deploy.yml",
            "Add runner jobs: ssh login, git pull, activate venv, pip install, migrate database, collect static, and pm2 restart ecosystem.config.js.",
            "Verify end-to-end automation triggers successfully on git push."
        ]
    )

    output_path = os.path.join("backend-django", "backend_deployment.pptx")
    prs.save(output_path)
    print(f"Generated: {output_path}")

if __name__ == "__main__":
    # Create output directories if they don't exist
    os.makedirs("forntend-next.js", exist_ok=True)
    os.makedirs("backend-django", exist_ok=True)
    
    generate_frontend()
    generate_backend()
